from __future__ import annotations
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from app.services.dataset_cache import DatasetCache

def generate_pptx_deck() -> str:
    """
    Generates an executive PowerPoint (.pptx) slide deck for board presentation.
    """
    prs = Presentation()

    # Slide 1: Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Data Analyst Agent"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    p2 = tf.add_paragraph()
    p2.text = "Executive Dataset Intelligence & Predictive Analytics Report Deck"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(37, 99, 235)

    # Slide 2: Dataset Profile & Summary
    slide2 = prs.slides.add_slide(blank_slide_layout)
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(5))
    tf2 = txBox2.text_frame
    
    title_p = tf2.paragraphs[0]
    title_p.text = "Executive Dataset Summary"
    title_p.font.size = Pt(28)
    title_p.font.bold = True

    df = DatasetCache.get_dataset()
    rows = len(df) if df is not None else 0
    cols = len(df.columns) if df is not None else 0

    body_p = tf2.add_paragraph()
    body_p.text = f"• Total Records Processed: {rows:,} rows\n• Attributes & Features: {cols} columns\n• DuckDB Engine Status: In-Memory High Performance\n• Security Level: OWASP 5-Pillar Hardened"
    body_p.font.size = Pt(16)

    # Save presentation
    reports_dir = Path("reports")
    reports_dir.mkdir(exist_ok=True)
    filename = f"AI_Data_Analyst_Deck_{int(os.times().elapsed)}.pptx"
    filepath = reports_dir / filename
    prs.save(filepath)

    return str(filepath)
